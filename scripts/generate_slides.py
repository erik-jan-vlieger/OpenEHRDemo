#!/usr/bin/env python3
"""
scripts/generate_slides.py
--------------------------
Genereert een professionele 16:9 PowerPoint-presentatie (.pptx) voor de Architectuurboard.
Deze kan direct in Google Drive geüpload worden en opent 1-op-1 als Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = REPO_ROOT / "docs"
DOCS_DIR.mkdir(parents=True, exist_ok=True)
PPTX_OUT = DOCS_DIR / "Sensire_openEHR_Architectuurboard_Presentatie.pptx"

# Sensire & openEHR Kleurenpalet
COLOR_PRIMARY = RGBColor(18, 52, 102)     # Diep marineblauw (#123466)
COLOR_BRAND = RGBColor(232, 131, 12)      # Sensire oranje (#e8830c)
COLOR_DARK = RGBColor(33, 37, 41)         # Donkergrijs/inkt (#212529)
COLOR_MUTED = RGBColor(108, 117, 125)     # Gedempt grijs (#6c757d)
COLOR_BG_LIGHT = RGBColor(248, 249, 250)  # Lichtgrijs paneel (#f8f9fa)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT_BG = RGBColor(238, 244, 255) # Lichtblauw (#eef4ff)
COLOR_CARD_BORDER = RGBColor(200, 215, 235)

def create_presentation():
    prs = Presentation()
    # 16:9 breedbeeldformaat (13.333 x 7.5 inch)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category="SENSIRE OPEN-EHR ARCHITECTUUR"):
        # Categorie banner
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_BRAND

        # Hoofdtitel
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY

        # Scheidingslijn
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BRAND
        line.line.color.rgb = COLOR_BRAND

    # ─────────────────────────────────────────────────────────────
    # SLIDE 1: TITELSLIDE
    # ─────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    # Oranje accent balk
    accent = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_BRAND
    accent.line.fill.background()

    # Titeltekst
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "SENSIRE DIGITALE INNOVATIE · ARCHITECTUURBOARD"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_BRAND
    p0.space_after = Pt(12)

    p1 = tf.add_paragraph()
    p1.text = "openEHR Ecosysteem & Runtime Architectuur"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Semantische Fabriek, Contract-First Decoupling en Runtime Roadmap naar EHRbase CDR"
    p2.font.size = Pt(17)
    p2.font.color.rgb = RGBColor(210, 225, 245)
    p2.space_after = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "Auteur: Digitale Innovatie & openEHR Werkgroep  |  Status: Definitief Architectuurvoorstel"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(160, 185, 215)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 2: EXECUTIVE SUMMARY & KERNPRINCIPES
    # ─────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Executive Summary: Drie Pijlers voor een Duurzaam Gezondheidsplatform")

    cards_data = [
        ("1. Data & App Scheiding", "Klinische gegevens worden opgeslagen in een leveranciersonafhankelijk Clinical Data Repository (openEHR CDR). Geen vendor lock-in in applicatiesilo's."),
        ("2. Contract-First & Decoupled", "Semantiek (archetypes, OPT 1.4 XML) wordt autonoom gecompileerd. Applicaties consumeren zuivere, snelle JSON-contracten zonder runtime Java-ballast."),
        ("3. Zero-Burden voor de Zorg", "Zowel wijkverpleegkundigen als zorgpad-redacteuren behouden hun vertrouwde functionele interface. De openEHR-transformatie verloopt 100% automatisch.")
    ]

    for i, (title, desc) in enumerate(cards_data):
        left = Inches(0.8 + i * 4.0)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(14)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 3: HET 3-LAGEN ARCHITECTUURMODEL
    # ─────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Het 3-Lagen Ecosysteem: Scheiding van Inhoud, Semantiek & Presentatie")

    layers = [
        ("LAAG 1: AUTEURS & KLINISCHE BRONNEN", "• SensireSpecieleBeslisOndersteuning (164 zorgpaden, ReAble flows, actieve dagelijkse redactie)\n• Sensii Wijkverpleegkundig EPD (58 surveys, Omaha-systematiek, SFP-beslismotor)"),
        ("LAAG 2: DE SEMANTISCHE FABRIEK (OpenEHRDemo)", "• 730+ CKM & 7 Custom Nederlandse Archetypes (Gordon, GFI, Omaha, Regietrede, VAG)\n• Compilatiestraat: Archie Java Flattener + full_opt14_lxml.py (299 gevalideerde OPTs)\n• Exporteert contracten: sensii_openehr_contract.json & carepaths_openehr_contract.json"),
        ("LAAG 3: PRESENTATIE & RUNTIME PERSISTENTIE", "• Live openEHR Inspector & Data-Twin in Sensii Doorloop\n• Lens 6: openEHR Blueprint & AQL Query Lab in Meta-Visualisaties\n• Runtime FLAT JSON Mapper & REST client naar EHRbase CDR (/ehr/{ehr_id}/composition)")
    ]

    for i, (l_title, l_desc) in enumerate(layers):
        top = Inches(1.8 + i * 1.7)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_ACCENT_BG if i == 1 else COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_BRAND if i == 1 else COLOR_CARD_BORDER

        tb = s3.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = l_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = l_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 4: DE 7 CUSTOM NEDERLANDSE ARCHETYPES
    # ─────────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "De 7 Custom Archetypes voor de Nederlandse Wijkverpleging")

    arch_list = [
        ("Groningen Frailty Indicator (GFI)", "OBSERVATION", "15 items over 4 domeinen (fysiek, cognitief, sociaal, psychisch), somscore 0-15, cutoff >= 4 (SNOMED 225544003)."),
        ("Regietrede & Zelfredzaamheid", "OBSERVATION", "Inschaling van 5 eigen regietreden met vergelijking motoradvies vs professional vs cliënt."),
        ("Omaha System Probleem & Zorgplan", "EVALUATION", "42 aandachtsgebieden over 4 domeinen, 4 actiesoorten, 76 actievlakken en KBS-streefscores (Kennis, Gedrag, Status)."),
        ("11 Gezondheidspatronen van Gordon", "OBSERVATION", "Gestructureerde functionele anamnese volgens de 11 verpleegkundige patronen van Marjory Gordon."),
        ("Wijkverpleegkundige Triage", "EVALUATION", "Triagebesluit passendheid wijkverpleging, urgentiecategorie en relatietype (Langer Thuis)."),
        ("Verpleegkundig Adviesgesprek (VAG)", "EVALUATION", "Gestructureerde vastlegging van de hulpvraag, wenselijke situatie, persoonlijke kracht en mantelzorgbalans."),
        ("Zorgfinanciering & Leveringsvorm", "ADMIN_ENTRY", "Wettelijk kader (Zvw, Wlz, Wmo), leveringsvorm (ZIN, PGB, VPT, MPT) en geïndiceerde zorgminuten.")
    ]

    for i, (a_name, a_rm, a_desc) in enumerate(arch_list):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.8 + row * 1.3)

        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.7), Inches(1.18))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08), Inches(5.4), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{a_name}  [{a_rm}]"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(3)

        p2 = tf.add_paragraph()
        p2.text = a_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 5: COMPILATIESTRAAT & TERMINOLOGIECASCADE
    # ─────────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Compilatiestraat (Archie Java & OPT 1.4) en 5-Staps Terminologiecascade")

    # Linker kolom: Compilatiestraat
    card_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_BG_LIGHT
    card_l.line.color.rgb = COLOR_CARD_BORDER

    tb_l = s5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "⚙️ Geautomatiseerde Compilatiestraat"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf_l.add_paragraph()
    p2.text = "1. Bron ADLT Templates (templates/*.adlt)\n   Definieert RM-structuur, slots en verplichte velden.\n\n2. Nedap Archie Java Flattener\n   Valideert RM-regels, expandeert slots en genereert OPT.\n\n3. Python full_opt14_lxml.py Postprocessor\n   Injecteert xsi:type, corrigeert namespaces en valideert XML.\n\n4. Gevalideerde OPT 1.4 XML (opts/*.opt)\n   299 gecompileerde bestanden gereed voor EHRbase CDR."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # Rechter kolom: Terminologiecascade
    card_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_ACCENT_BG
    card_r.line.color.rgb = COLOR_BRAND

    tb_r = s5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "📚 5-Staps Terminologiecascade"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf_r.add_paragraph()
    p2.text = "Strikte volgorde voor semantische binding:\n\n1. Nictiz ZIBs (Zorginformatiebouwstenen)\n   Primaire bron voor landelijke interoperabiliteit.\n\n2. SNOMED CT NL Editie\n   Internationale en Nederlandse klinische terminologie.\n\n3. openEHR CKM Terminology\n   Officiële archetype internal value-sets en bindings.\n\n4. V&VN Richtlijnen & Omaha System NL\n   Beroepsstandaarden voor de wijkverpleging.\n\n5. Zorgpad Fallback\n   Sensire-specifieke terminologie indien geen standaard."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 6: CONTRACT-FIRST DATA DISTRIBUTIE
    # ─────────────────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Contract-First Data Distributie: Waarom Sensii en Zorgpaden Schoon Blijven")

    card1 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_BG_LIGHT
    card1.line.color.rgb = COLOR_CARD_BORDER

    tb1 = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "📦 sensii_openehr_contract.json (39.5 KB)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "• Mapt alle 58 surveys op de 5 templates en 7 custom archetypes.\n• Bevat exacte RM node-paden (at00xx codes) voor elk formulier.\n• Bevat kant-en-klare AQL sample queries voor data-ontsluiting.\n• Bevat dynamische triggers tussen Omaha aandachtsgebieden en de 164 speciële zorgpaden (bijv. Huid ➔ Diabetische Voet).\n• Bevat SHA256 integriteits-checksum."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    card2 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_BG_LIGHT
    card2.line.color.rgb = COLOR_CARD_BORDER

    tb2 = s6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "📦 carepaths_openehr_contract.json (293 paden)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf2.add_paragraph()
    p2.text = "• Bevat het register van alle 164 zorgpaden (regulier & ReAble).\n• Mapt elk zorgpad op zijn gecompileerde OPT 1.4 XML bestand.\n• Definieert de ReAble-status (inclusief ACTION.health_education en doel-tracking).\n• Dient als basis voor de niet-verstorende Fase 2 openEHR badge in de WebViewer."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 7: VISUALISATIES & DUAL-PERSPECTIEF
    # ─────────────────────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Visualisatie & Inspectie: Het Dual-Perspectief in de Praktijk")

    card_v1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_v1.fill.solid()
    card_v1.fill.fore_color.rgb = COLOR_BG_LIGHT
    card_v1.line.color.rgb = COLOR_CARD_BORDER

    tb_v1 = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_v1 = tb_v1.text_frame
    tf_v1.word_wrap = True
    p = tf_v1.paragraphs[0]
    p.text = "👩‍⚕️ De Doorloop: Live openEHR Inspector"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf_v1.add_paragraph()
    p2.text = "• Klinische Weergave (Wijkverpleegkundige):\n  Rustig en herkenbaar cliëntdossier (Gordon, Omaha KBS, GFI, regietrede, indicatie).\n\n• openEHR Inspector (Technicus & Architect):\n  Live data-twin met RM-klasse, archetype ID, live FLAT JSON payload fragment en de bijbehorende AQL-query.\n\n• Live Sync Pulse: Directe feedback over compositie-validatie."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    card_v2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_v2.fill.solid()
    card_v2.fill.fore_color.rgb = COLOR_BG_LIGHT
    card_v2.line.color.rgb = COLOR_CARD_BORDER

    tb_v2 = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
    tf_v2 = tb_v2.text_frame
    tf_v2.word_wrap = True
    p = tf_v2.paragraphs[0]
    p.text = "🧬 Meta-Visualisaties: Lens 6 Blueprint & Query Lab"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    p2 = tf_v2.add_paragraph()
    p2.text = "• Dekkingsgraad Dashboard:\n  100% dekking over 58 formulieren en 42 Omaha-gebieden.\n\n• Interactieve Template-Hiërarchie:\n  Visuele compositieboom (Mermaid/D3).\n\n• Live AQL Query Tester:\n  Interactieve query-editor waar analisten en IT direct queries kunnen uitvoeren op voorbeeld-dossiers."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 8: RUNTIME ROADMAP: SCHRIJVEN NAAR EHRBASE
    # ─────────────────────────────────────────────────────────────
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Runtime Roadmap: 4 Stappen om Live naar EHRbase CDR te Schrijven")

    steps = [
        ("1. EHRbase Server & OPT Upload", "EHRbase instance (Docker/K8s + PostgreSQL) inrichten. De 5 OPT 1.4 XML bestanden eenmalig registreren via POST /rest/openehr/v1/definition/template/adl1.4."),
        ("2. WebTemplate FLAT JSON Mapper", "Lichtgewicht Python helper (sensii_openehr_mapper.py) in Sensii backend die formulier-antwoorden omzet naar Simplified FLAT JSON formaat."),
        ("3. EHRbase REST API Client", "HTTP POST client die de payload verzendt naar /ehr/{ehr_id}/composition met template-id header en OAuth2 Bearer token."),
        ("4. Audit, Versioning & AQL", "EHRbase valideert tegen de OPT, slaat onveranderlijk op met ObjectVersionID en stelt data direct beschikbaar voor realtime AQL bevraging.")
    ]

    for i, (st_title, st_desc) in enumerate(steps):
        top = Inches(1.8 + i * 1.3)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s8.shapes.add_textbox(Inches(1.0), top + Inches(0.08), Inches(11.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = st_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(3)

        p2 = tf.add_paragraph()
        p2.text = st_desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 9: ADOPTIESTRATEGIE VOOR SPECIËLE BESLISONDERSTEUNING
    # ─────────────────────────────────────────────────────────────
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Adoptiestrategie Speciële Beslisondersteuning: Fase 2 Passieve Schaduw-Adapter")

    fases = [
        ("FASE 1: ZERO-TOUCH (Actueel)", "• Redacteurs blijven 100% werken zoals nu (Markdown/Mermaids/JSON).\n• OpenEHRDemo genereert de openEHR templates volledig autonoom op de achtergrond.\n• Geen enkele verstoring van het dagelijkse zorgpaden-onderhoud."),
        ("FASE 2: PASSIEVE INFORMATIEVE BADGE (Aanbevolen Stap)", "• In de WebViewer verschijnt een subtiele badge: [ 🧬 openEHR Template Beschikbaar ].\n• Klikken opent een inspectie-popup met template-ID, archetypes en AQL snippet.\n• Volledig optioneel en niet-blokkerend voor redacteurs."),
        ("FASE 3: OPT-IN AI ASSISTENT (Toekomst)", "• Vrijwillige AI-assistent die bij het toevoegen van een nieuwe zorgpadstap suggesties doet voor Omaha-actievlakken of SNOMED-termen.\n• Altijd adviserend, nooit dwingend.")
    ]

    for i, (f_title, f_desc) in enumerate(fases):
        top = Inches(1.8 + i * 1.7)
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_ACCENT_BG if i == 1 else COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_BRAND if i == 1 else COLOR_CARD_BORDER

        tb = s9.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = f_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 10: CONCLUSIE & VOORGESTELDE BESLUITEN
    # ─────────────────────────────────────────────────────────────
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Conclusie & Voorgestelde Besluiten voor het Architectuurboard")

    decisions = [
        ("1. Goedkeuring Ecosysteem Architectuur", "Het 3-lagenmodel met contract-first scheiding tussen Semantische Fabriek (OpenEHRDemo) en consumenten-applicaties (Sensii en Speciële Beslisondersteuning) vaststellen als Sensire enterprise standaard."),
        ("2. Goedkeuring Fase 2 in Speciële Beslisondersteuning", "Inrichten van de passieve, niet-verstorende openEHR inspectie-badge in de WebViewer conform het geleverde JSON-contract."),
        ("3. Goedkeuring Runtime Pilot EHRbase", "Starten van een gecontroleerde pilot om de Sensii Doorloop via FLAT JSON composities live weg te schrijven naar een dedicated EHRbase CDR instance.")
    ]

    for i, (d_title, d_desc) in enumerate(decisions):
        top = Inches(1.8 + i * 1.7)
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = s10.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.3), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = d_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = d_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_DARK

    # Opslaan
    prs.save(PPTX_OUT)
    print(f"✅ PowerPoint presentatie succesvol gegenereerd: {PPTX_OUT} ({PPTX_OUT.stat().st_size / 1024:.1f} KB)")
    return PPTX_OUT

if __name__ == "__main__":
    create_presentation()
